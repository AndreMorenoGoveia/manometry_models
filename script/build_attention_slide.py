"""Build a presentation slide showing the graph attention for one image per class.

Renders a single composite panel (input bands on top, effective contribution
below) and appends it to a copy of the deck, matching the deck's visual style.
"""

from __future__ import annotations

import argparse
import shutil
from pathlib import Path

BACKGROUND = "F6F4F0"
ACCENT = "E2661A"
EYEBROW = "1F8F9C"
TITLE = "142A3E"
BODY = "33424F"
CAPTION = "6B7C8C"


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Render the graph-attention slide for the Wang CVP-GAT deck.")
    parser.add_argument("--checkpoint", type=Path, required=True, help="Path to best_model.pt.")
    parser.add_argument(
        "--test-dir",
        type=Path,
        required=True,
        help="ImageFolder-style directory; one image is taken from each class subfolder.",
    )
    parser.add_argument("--deck", type=Path, required=True, help="Existing .pptx to copy and extend.")
    parser.add_argument("--output-deck", type=Path, required=True, help="Where to write the extended deck.")
    parser.add_argument(
        "--panel-path",
        type=Path,
        default=Path("reports/graph_attention/painel_classes.png"),
        help="Where to write the composite panel image.",
    )
    parser.add_argument("--device", type=str, default="cpu", choices=["auto", "cpu", "cuda", "mps"])
    return parser


def main() -> None:
    args = build_parser().parse_args()

    import matplotlib

    matplotlib.use("Agg")

    import matplotlib.pyplot as plt
    import numpy as np
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    from manometry_models.graph_attention import GraphAttentionExplainer
    from manometry_models.training import resolve_device

    if not args.deck.exists():
        raise FileNotFoundError(f"Deck not found: {args.deck}")

    explainer = GraphAttentionExplainer(args.checkpoint, resolve_device(args.device))

    class_dirs = sorted(path for path in args.test_dir.iterdir() if path.is_dir())
    if not class_dirs:
        raise ValueError(f"No class subfolders found in {args.test_dir}")

    results = []
    for class_dir in class_dirs:
        images = sorted(path for path in class_dir.iterdir() if path.suffix.lower() in {".jpg", ".jpeg", ".png"})
        if not images:
            continue
        result = explainer.explain(images[0])
        results.append((class_dir.name, result))

    columns = len(results)
    figure, axes = plt.subplots(2, columns, figsize=(2.3 * columns, 5.0), squeeze=False)
    figure.patch.set_facecolor("#" + BACKGROUND)

    for column, (class_name, result) in enumerate(results):
        top = axes[0][column]
        top.imshow(result.display)
        for node in range(len(result.band_starts)):
            top.axhline(result.band_starts[node] * result.cell, color="white", linewidth=0.7, alpha=0.8)
            top.text(
                3,
                result.band_center(node),
                f"n{node}",
                color="white",
                fontsize=6,
                va="center",
                bbox={"facecolor": "black", "alpha": 0.4, "pad": 0.8, "edgecolor": "none"},
            )
        correct = "" if result.predicted_class == class_name else f"\n(previu {result.predicted_class})"
        top.set_title(f"{class_name}{correct}", fontsize=10, color="#" + TITLE, fontweight="bold", pad=6)
        top.set_xticks([])
        top.set_yticks([])

        bottom = axes[1][column]
        bottom.imshow(result.display)
        # The HRM plot is already a rainbow colormap, so a second colormap on
        # top is unreadable. Dim the unused bands instead and let the original
        # colors survive where the graph actually looked.
        profile = result.band_profile(result.contribution)
        spread = profile.max() - profile.min()
        normalized = (profile - profile.min()) / spread if spread > 0 else np.ones_like(profile)
        veil = np.zeros((result.image_size, result.image_size, 4))
        veil[..., 3] = np.repeat((0.78 * (1.0 - normalized))[:, None], result.image_size, axis=1)
        bottom.imshow(veil)
        peak = int(np.argmax(result.contribution))
        for node, weight in enumerate(result.contribution):
            bottom.text(
                result.image_size - 4,
                result.band_center(node),
                f"{weight:.2f}",
                color="white",
                fontsize=6.5,
                ha="right",
                va="center",
                fontweight="bold" if node == peak else "normal",
                bbox={
                    "facecolor": "#" + ACCENT if node == peak else "black",
                    "alpha": 0.65,
                    "pad": 0.8,
                    "edgecolor": "none",
                },
            )
        bottom.set_xticks([])
        bottom.set_yticks([])

    axes[0][0].set_ylabel("entrada 224²\n6 nós", fontsize=8.5, color="#" + BODY)
    axes[1][0].set_ylabel("contribuição efetiva\n(rollout × readout)", fontsize=8.5, color="#" + BODY)
    figure.subplots_adjust(left=0.055, right=0.995, top=0.90, bottom=0.01, wspace=0.06, hspace=0.06)

    args.panel_path.parent.mkdir(parents=True, exist_ok=True)
    figure.savefig(args.panel_path, dpi=200, facecolor=figure.get_facecolor())
    plt.close(figure)

    presentation = Presentation(str(args.deck))
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    def rectangle(left, top, width, height, color) -> None:
        from pptx.enum.shapes import MSO_SHAPE

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color)
        shape.line.fill.background()
        shape.shadow.inherit = False

    rectangle(0, 0, presentation.slide_width, presentation.slide_height, BACKGROUND)
    rectangle(0, 0, presentation.slide_width, Inches(0.12), ACCENT)

    def textbox(left, top, width, height):
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        return frame

    frame = textbox(Inches(0.70), Inches(0.42), Inches(11.90), Inches(0.30))
    run = frame.paragraphs[0].add_run()
    run.text = "ATENÇÃO DO GRAFO · UM EXEMPLO POR CLASSE"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.name = "Consolas"
    run.font.color.rgb = RGBColor.from_string(EYEBROW)

    frame = textbox(Inches(0.70), Inches(0.72), Inches(11.90), Inches(0.90))
    run = frame.paragraphs[0].add_run()
    run.text = "O que o grafo realmente olha, por diagnóstico"
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor.from_string(TITLE)

    slide.shapes.add_picture(str(args.panel_path), Inches(0.70), Inches(1.68), width=Inches(11.90))

    frame = textbox(Inches(0.70), Inches(6.30), Inches(11.90), Inches(0.95))
    segments = [
        ("Linha de cima: ", True, BODY),
        ("as 6 bandas em que a imagem é dividida (sobrepostas).  ", False, BODY),
        ("Linha de baixo: ", True, BODY),
        ("peso de cada região no vetor que decide a classe.\n", False, BODY),
        ("Atenção: ", True, ACCENT),
        (
            "o readout aprendido colapsou — escolhe o nó 2 em 90,5% do conjunto de teste "
            "(entropia 0,035 nats vs. 1,792 do uniforme). O que varia entre diagnósticos é o "
            "rollout do GAT, ou seja, quais regiões alimentam esse nó.",
            False,
            BODY,
        ),
    ]
    paragraph = frame.paragraphs[0]
    for text, bold, color in segments:
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(12)
        run.font.bold = bold
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor.from_string(color)

    frame = textbox(Inches(0.70), Inches(7.16), Inches(11.90), Inches(0.26))
    run = frame.paragraphs[0].add_run()
    run.text = (
        f"Checkpoint: {args.checkpoint}  ·  imagens de {args.test_dir}  ·  "
        "rollout com residual 50/50 (heurística de Abnar & Zuidema)"
    )
    run.font.size = Pt(9)
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor.from_string(CAPTION)

    if args.output_deck.resolve() != args.deck.resolve():
        shutil.copyfile(args.deck, args.output_deck)
    presentation.save(str(args.output_deck))

    print(f"Painel: {args.panel_path}")
    print(f"Slide {len(presentation.slides)} adicionado em: {args.output_deck}")
    for class_name, result in results:
        weights = ", ".join(f"n{i}={w:.2f}" for i, w in enumerate(result.contribution))
        print(f"  {class_name:20s} -> {result.predicted_class:20s} ({result.probabilities.max():.3f})  {weights}")


if __name__ == "__main__":
    main()
